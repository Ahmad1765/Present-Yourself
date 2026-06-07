"""Render a SlideBlueprint into a PPTX byte stream using python-pptx.

Minimum viable renderer: supports title, section_header, bullet_list, two_column,
image_caption, quote, stat_callout, comparison_table, chart_placeholder, closing.
Uses the deck's design_tokens for palette + fonts; logos are placed when present.
"""

from __future__ import annotations

import io
import re
from urllib.request import urlopen

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from app.schemas.slide import SlideBlueprint


def render_blueprint(bp: SlideBlueprint) -> bytes:
    deck = bp.deck
    prs = Presentation()
    prs.slide_width = Emu(deck.design_tokens.dimensions.width_emu)
    prs.slide_height = Emu(deck.design_tokens.dimensions.height_emu)

    palette = deck.design_tokens.palette
    fonts = deck.design_tokens.fonts
    bg = _hex(palette.background)
    fg = _hex(palette.foreground)
    accent = _hex(palette.accent1)
    muted = _hex(palette.muted)

    blank = prs.slide_layouts[6]  # Blank
    for slide_spec in deck.slides:
        slide = prs.slides.add_slide(blank)
        _set_bg(slide, bg)
        title = _find_text(slide_spec, role="title")
        subtitle = _find_text(slide_spec, role="subtitle")
        body = _find_text(slide_spec, role="body")
        bullets = _find_bullets(slide_spec)
        image_elem = _find_image(slide_spec)
        caption = _find_text(slide_spec, role="caption")

        layout = slide_spec.layout_hint or _default_layout(slide_spec.type)

        if slide_spec.type in ("title", "closing"):
            _add_centered_title(slide, title or deck.title, accent, fonts.heading.family)
            if subtitle:
                _add_centered_subtitle(slide, subtitle, muted, fonts.body.family)
        elif slide_spec.type == "section_header":
            _add_centered_title(slide, title or "Section", accent, fonts.heading.family)
        elif slide_spec.type in ("bullet_list",):
            _add_title(slide, title, fg, fonts.heading.family)
            _add_bullets(slide, bullets, fg, fonts.body.family)
        elif slide_spec.type == "two_column":
            _add_title(slide, title, fg, fonts.heading.family)
            cols = _split_two_columns(slide_spec)
            _add_bullets(slide, cols[0], fg, fonts.body.family, left=Inches(0.5), width=Inches(5.8))
            _add_bullets(slide, cols[1], fg, fonts.body.family, left=Inches(6.5), width=Inches(5.8))
        elif slide_spec.type == "image_caption":
            _add_title(slide, title, fg, fonts.heading.family)
            _add_image(slide, image_elem, layout)
            if caption:
                _add_caption(slide, caption, muted, fonts.body.family)
        elif slide_spec.type == "quote":
            _add_quote(slide, title or body or "", accent, fonts.heading.family)
        elif slide_spec.type == "stat_callout":
            _add_stat(slide, title or "", subtitle or "", accent, muted, fonts.heading.family)
        elif slide_spec.type == "comparison_table":
            _add_title(slide, title, fg, fonts.heading.family)
            _add_caption(slide, body or "Comparison", muted, fonts.body.family)
        elif slide_spec.type == "chart_placeholder":
            _add_title(slide, title, fg, fonts.heading.family)
            _add_placeholder(slide, "Chart placeholder", muted, fonts.body.family)
        else:
            _add_title(slide, title or slide_spec.type, fg, fonts.heading.family)

        if slide_spec.notes:
            slide.notes_slide.notes_text_frame.text = slide_spec.notes

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


# ─── helpers ─────────────────────────────────────────────────────────────────

_HEX_RE = re.compile(r"^#?([0-9a-fA-F]{6})$")


def _hex(value: str) -> RGBColor:
    m = _HEX_RE.match(value or "")
    if not m:
        return RGBColor(0x0F, 0x17, 0x2A)
    raw = m.group(1)
    return RGBColor(int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16))


def _set_bg(slide, color: RGBColor) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _find_text(slide_spec, *, role: str) -> str | None:
    for el in slide_spec.elements:
        if el.kind == "text" and el.role == role and el.content:
            return el.content
    return None


def _find_bullets(slide_spec) -> list[str]:
    for el in slide_spec.elements:
        if el.kind == "bullets" and el.items:
            return el.items
    return []


def _find_image(slide_spec):
    for el in slide_spec.elements:
        if el.kind == "image":
            return el
    return None


def _split_two_columns(slide_spec) -> tuple[list[str], list[str]]:
    bullets = _find_bullets(slide_spec)
    half = max(1, len(bullets) // 2)
    return bullets[:half], bullets[half:]


def _default_layout(slide_type: str) -> str:
    return {
        "title": "centered",
        "section_header": "centered",
        "bullet_list": "left_text_only",
        "two_column": "two_column",
        "image_caption": "right_image",
        "quote": "centered",
        "stat_callout": "centered",
        "comparison_table": "table",
        "chart_placeholder": "left_text_only",
        "closing": "centered",
    }.get(slide_type, "left_text_only")


def _add_title(slide, text: str | None, color: RGBColor, font_family: str) -> None:
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.0), Inches(1.0))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text or ""
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.name = font_family
    run.font.color.rgb = color


def _add_centered_title(slide, text: str, color: RGBColor, font_family: str) -> None:
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.0), Inches(1.6))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    run = p.runs[0]
    run.font.size = Pt(54)
    run.font.bold = True
    run.font.name = font_family
    run.font.color.rgb = color


def _add_centered_subtitle(slide, text: str, color: RGBColor, font_family: str) -> None:
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(4.1), Inches(12.0), Inches(0.8))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    run = p.runs[0]
    run.font.size = Pt(22)
    run.font.name = font_family
    run.font.color.rgb = color


def _add_bullets(
    slide,
    items: list[str],
    color: RGBColor,
    font_family: str,
    *,
    left=Inches(0.5),
    top=Inches(1.6),
    width=Inches(12.0),
    height=Inches(5.5),
) -> None:
    if not items:
        return
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {item}"
        run = p.runs[0]
        run.font.size = Pt(20)
        run.font.name = font_family
        run.font.color.rgb = color
        p.space_after = Pt(8)


def _add_image(slide, image_elem, layout: str) -> None:
    if image_elem is None or not image_elem.source or not image_elem.source.get("url"):
        return
    try:
        with urlopen(image_elem.source["url"], timeout=8) as resp:
            data = resp.read()
        stream = io.BytesIO(data)
        if "right_image" in layout:
            slide.shapes.add_picture(stream, Inches(7.0), Inches(1.6), width=Inches(5.8), height=Inches(4.5))
        elif "left_image" in layout:
            slide.shapes.add_picture(stream, Inches(0.5), Inches(1.6), width=Inches(5.8), height=Inches(4.5))
        else:
            slide.shapes.add_picture(stream, Inches(3.0), Inches(2.0), width=Inches(7.3), height=Inches(4.0))
    except Exception:  # noqa: BLE001
        pass


def _add_caption(slide, text: str, color: RGBColor, font_family: str) -> None:
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(12.0), Inches(0.6))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = text
    run = p.runs[0]
    run.font.size = Pt(14)
    run.font.italic = True
    run.font.name = font_family
    run.font.color.rgb = color


def _add_quote(slide, text: str, color: RGBColor, font_family: str) -> None:
    tx = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.0), Inches(2.5))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = f"“{text}”"
    run = p.runs[0]
    run.font.size = Pt(40)
    run.font.italic = True
    run.font.name = font_family
    run.font.color.rgb = color


def _add_stat(slide, value: str, label: str, value_color: RGBColor, label_color: RGBColor, font_family: str) -> None:
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(2.4), Inches(12.0), Inches(2.0))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = value
    run = p.runs[0]
    run.font.size = Pt(96)
    run.font.bold = True
    run.font.name = font_family
    run.font.color.rgb = value_color
    tx2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(12.0), Inches(0.8))
    tf2 = tx2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    p2.text = label
    run2 = p2.runs[0]
    run2.font.size = Pt(22)
    run2.font.name = font_family
    run2.font.color.rgb = label_color


def _add_placeholder(slide, text: str, color: RGBColor, font_family: str) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.0), Inches(2.0), Inches(9.0), Inches(4.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF4)
    shape.line.color.rgb = RGBColor(0xE7, 0xE5, 0xE4)
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    run = p.runs[0]
    run.font.size = Pt(20)
    run.font.name = font_family
    run.font.color.rgb = color
